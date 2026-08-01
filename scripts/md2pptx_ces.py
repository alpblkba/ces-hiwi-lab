#!/usr/bin/env python3
"""
Converts the Marp lab decks into PowerPoint files styled in the KIT/CES colours.

This is a stand-in for building the decks on the host machine with
`md2pptx.py --template CES_ppt_template-1.pptx`, which is the route that
applies the official corporate template. Use this one to review content and
layout on the server, where the template file is not available.

    python3 scripts/md2pptx_ces.py labs/lab1/docs/lab1-hls-bringup.md
    python3 scripts/md2pptx_ces.py labs/lab1/docs/lab1-hls-bringup.md -o out.pptx

Supported Marp constructs: headings, bullet lists, fenced code blocks,
pipe tables, and the deck's <div class="shot|flow|note|warn"> blocks.
"""

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

KIT_GREEN = RGBColor(0x00, 0x96, 0x82)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x66, 0x66, 0x66)
CODE_BG = RGBColor(0xF2, 0xF4, 0xF3)
NOTE_BG = RGBColor(0xE8, 0xF4, 0xF2)
WARN_BG = RGBColor(0xFD, 0xF0, 0xE8)
WARN_LINE = RGBColor(0xD0, 0x5A, 0x1E)
SHOT_BG = RGBColor(0xF7, 0xF9, 0xF8)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)
BODY_TOP = Inches(1.42)
BODY_W = SLIDE_W - 2 * MARGIN


def strip_inline(text, keep_indent=False):
    """Removes Marp/HTML inline markup, keeping the readable text.

    keep_indent preserves leading whitespace, which matters for the <div
    class="flow"> blocks: they are column-aligned ASCII diagrams and lose
    their meaning if the indentation is collapsed.
    """
    text = re.sub(r"<br\s*/?>", " ", text)
    text = re.sub(r"<[^>]+>", "", text)
    text = text.replace("&nbsp;", " ").replace("&amp;", "&")
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"!\[[^\]]*\]\([^)]*\)", "", text)
    return text.rstrip() if keep_indent else text.strip()


def parse_slides(md_text):
    """Splits the Marp source into slides, dropping frontmatter and comments."""
    md_text = re.sub(r"^---\n.*?\n---\n", "", md_text, count=1, flags=re.S)
    md_text = re.sub(r"<!--.*?-->", "", md_text, flags=re.S)
    raw = re.split(r"^---\s*$", md_text, flags=re.M)
    return [s for s in (b.strip() for b in raw) if s]


def parse_blocks(slide_md):
    """Turns one slide's markdown into an ordered list of typed blocks."""
    blocks, lines, i = [], slide_md.split("\n"), 0
    bullets, table = [], []

    def flush():
        nonlocal bullets, table
        if bullets:
            blocks.append(("bullets", bullets))
            bullets = []
        if table:
            blocks.append(("table", table))
            table = []

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        m = re.match(r"^(#{1,3})\s+(.*)", stripped)
        if m:
            flush()
            blocks.append(("h%d" % len(m.group(1)), strip_inline(m.group(2))))
            i += 1
            continue

        if stripped.startswith("```"):
            flush()
            i += 1
            code = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1
            blocks.append(("code", "\n".join(code)))
            continue

        m = re.match(r'^<div class="(shot|shot-row|flow|note|warn)[^"]*">', stripped)
        if m:
            flush()
            kind = m.group(1)
            chunk = []
            while i < len(lines) and "</div>" not in lines[i]:
                chunk.append(lines[i])
                i += 1
            if i < len(lines):
                chunk.append(lines[i])
            i += 1
            body = "\n".join(chunk)
            body = re.sub(r'^<div class="[^"]*">', "", body.strip())
            body = body.replace("</div>", "")
            if kind.startswith("shot"):
                name = re.search(r"img/[\w\-.]+", body)
                desc = strip_inline(re.sub(r"img/[\w\-.]+", "", body))
                blocks.append(("shot", (name.group(0) if name else "screenshot", desc)))
            else:
                keep = kind == "flow"   # flow blocks are column-aligned diagrams
                text = "\n".join(
                    strip_inline(x, keep_indent=keep) for x in body.split("\n"))
                blocks.append((kind, text.strip("\n") if keep else text.strip()))
            continue

        if stripped.startswith("|"):
            cells = [c.strip() for c in stripped.strip("|").split("|")]
            if not all(re.fullmatch(r":?-{2,}:?", c) for c in cells):
                table.append([strip_inline(c) for c in cells])
            i += 1
            continue

        if re.match(r"^[-*]\s+", stripped):
            if table:
                flush()
            bullets.append(strip_inline(re.sub(r"^[-*]\s+", "", stripped)))
            i += 1
            continue

        if stripped:
            flush()
            blocks.append(("para", strip_inline(stripped)))
        i += 1

    flush()
    return blocks


def add_box(slide, top, height, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, top, BODY_W, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line:
        box.line.color.rgb = line
        box.line.width = Pt(1.25)
    else:
        box.line.fill.background()
    box.shadow.inherit = False
    box.adjustments[0] = 0.03
    return box


def set_text(frame, text, size, color=DARK, mono=False, bold=False, space=6):
    frame.word_wrap = True
    lines = text.split("\n")
    for idx, ln in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = ln
        p.space_after = Pt(space)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = "Consolas" if mono else "Arial"
    return frame


def build(md_path, out_path):
    slides_md = parse_slides(Path(md_path).read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    for n, sm in enumerate(slides_md):
        blocks = parse_blocks(sm)
        if not blocks:
            continue
        slide = prs.slides.add_slide(blank)
        is_title = n == 0

        if is_title:
            # The rule sits below the title block, which may wrap to two lines.
            band = slide.shapes.add_shape(1, 0, Inches(3.05), SLIDE_W, Inches(0.09))
            band.fill.solid()
            band.fill.fore_color.rgb = KIT_GREEN
            band.line.fill.background()
            band.shadow.inherit = False
            y = Inches(1.35)
            for kind, val in blocks:
                if kind == "h1":
                    tb = slide.shapes.add_textbox(MARGIN, y, BODY_W, Inches(1.6))
                    set_text(tb.text_frame, val, 38, DARK, bold=True)
                    y = Inches(3.35)
                elif kind == "h2":
                    tb = slide.shapes.add_textbox(MARGIN, y, BODY_W, Inches(0.7))
                    set_text(tb.text_frame, val, 21, KIT_GREEN)
                    y += Inches(0.78)
                elif kind == "para":
                    tb = slide.shapes.add_textbox(MARGIN, y, BODY_W, Inches(0.45))
                    set_text(tb.text_frame, val, 13, GREY)
                    y += Inches(0.42)
            continue

        y = Inches(0.5)
        title = next((v for k, v in blocks if k in ("h1", "h2")), None)
        if title:
            tb = slide.shapes.add_textbox(MARGIN, y, BODY_W, Inches(0.72))
            set_text(tb.text_frame, title, 27, DARK, bold=True)
            rule = slide.shapes.add_shape(1, MARGIN, Inches(1.22), Inches(1.7), Inches(0.045))
            rule.fill.solid()
            rule.fill.fore_color.rgb = KIT_GREEN
            rule.line.fill.background()
            rule.shadow.inherit = False
        y = BODY_TOP

        for kind, val in blocks:
            if kind in ("h1", "h2") and val == title:
                continue
            avail = SLIDE_H - Inches(0.5) - y
            if avail < Inches(0.5):
                break

            if kind == "bullets":
                h = min(Inches(0.34) * len(val) + Inches(0.12), avail)
                tb = slide.shapes.add_textbox(MARGIN, y, BODY_W, h)
                tf = tb.text_frame
                tf.word_wrap = True
                for i, b in enumerate(val):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = "▪  " + b
                    p.space_after = Pt(7)
                    for r in p.runs:
                        r.font.size = Pt(14)
                        r.font.color.rgb = DARK
                        r.font.name = "Arial"
                y += h + Inches(0.14)

            elif kind == "code":
                nl = len(val.split("\n"))
                h = min(Inches(0.215) * nl + Inches(0.26), avail)
                add_box(slide, y, h, CODE_BG)
                tb = slide.shapes.add_textbox(
                    MARGIN + Inches(0.16), y + Inches(0.1),
                    BODY_W - Inches(0.32), h - Inches(0.2))
                set_text(tb.text_frame, val, 10.5, DARK, mono=True, space=1)
                y += h + Inches(0.14)

            elif kind == "flow":
                nl = len(val.split("\n"))
                h = min(Inches(0.215) * nl + Inches(0.26), avail)
                add_box(slide, y, h, CODE_BG)
                tb = slide.shapes.add_textbox(
                    MARGIN + Inches(0.16), y + Inches(0.1),
                    BODY_W - Inches(0.32), h - Inches(0.2))
                set_text(tb.text_frame, val, 10.5, DARK, mono=True, space=1)
                y += h + Inches(0.14)

            elif kind in ("note", "warn"):
                nl = max(1, len(val) // 110 + 1)
                h = min(Inches(0.3) * nl + Inches(0.26), avail)
                add_box(slide, y, h,
                        NOTE_BG if kind == "note" else WARN_BG,
                        KIT_GREEN if kind == "note" else WARN_LINE)
                tb = slide.shapes.add_textbox(
                    MARGIN + Inches(0.18), y + Inches(0.09),
                    BODY_W - Inches(0.36), h - Inches(0.18))
                set_text(tb.text_frame, val, 12, DARK, space=2)
                y += h + Inches(0.14)

            elif kind == "shot":
                name, desc = val
                h = min(Inches(1.9), avail)
                box = add_box(slide, y, h, SHOT_BG, GREY)
                box.line.dash_style = 4
                tb = slide.shapes.add_textbox(
                    MARGIN + Inches(0.2), y + Inches(0.22),
                    BODY_W - Inches(0.4), h - Inches(0.4))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = name
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(13)
                    r.font.bold = True
                    r.font.color.rgb = KIT_GREEN
                    r.font.name = "Consolas"
                if desc:
                    p2 = tf.add_paragraph()
                    p2.text = desc
                    p2.alignment = PP_ALIGN.CENTER
                    for r in p2.runs:
                        r.font.size = Pt(11)
                        r.font.color.rgb = GREY
                        r.font.name = "Arial"
                y += h + Inches(0.14)

            elif kind == "table":
                rows, cols = len(val), max(len(r) for r in val)
                h = min(Inches(0.32) * rows, avail)
                shape = slide.shapes.add_table(rows, cols, MARGIN, y, BODY_W, h)
                tbl = shape.table
                # Suppress the default blue theme so the table matches KIT green.
                tbl.first_row = False
                tbl.horz_banding = False
                for ri, row in enumerate(val):
                    for ci in range(cols):
                        cell = tbl.cell(ri, ci)
                        cell.text = row[ci] if ci < len(row) else ""
                        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                        cell.fill.solid()
                        if ri == 0:
                            cell.fill.fore_color.rgb = KIT_GREEN
                        else:
                            cell.fill.fore_color.rgb = (
                                RGBColor(0xFF, 0xFF, 0xFF) if ri % 2
                                else RGBColor(0xF2, 0xF4, 0xF3))
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                r.font.size = Pt(11.5)
                                r.font.name = "Arial"
                                r.font.bold = ri == 0
                                r.font.color.rgb = (
                                    RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else DARK)
                y += h + Inches(0.14)

            elif kind == "para":
                h = min(Inches(0.3) * (len(val) // 110 + 1) + Inches(0.06), avail)
                tb = slide.shapes.add_textbox(MARGIN, y, BODY_W, h)
                set_text(tb.text_frame, val, 13, DARK)
                y += h + Inches(0.08)

        # page number
        pn = slide.shapes.add_textbox(
            SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.46), Inches(0.5), Inches(0.3))
        set_text(pn.text_frame, str(n), 10, GREY)

    prs.save(out_path)
    return len(prs.slides.__iter__.__self__._sldIdLst)


if __name__ == "__main__":
    args = [a for a in sys.argv[1:] if not a.startswith("-")]
    if not args:
        print(__doc__)
        sys.exit(1)
    src = Path(args[0])
    out = Path(sys.argv[sys.argv.index("-o") + 1]) if "-o" in sys.argv \
        else src.with_suffix(".pptx")
    n = build(src, out)
    print(f"wrote {out}  ({n} slides)")
